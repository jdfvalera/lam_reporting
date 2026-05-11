from io import BytesIO

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from charts.summary_charts import overview_chart
from charts.demographic_charts import demographic_chart
from charts.area_charts import area_chart
from charts.product_charts import product_pie_chart
from charts.creative_charts import creative_chart


# ── Slide constants ───────────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
HDR = Inches(0.72)
PAD = Inches(0.28)
RECT = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE


# ── Color helpers ─────────────────────────────────────────────────────────────
def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_RED  = RGBColor(0xFC, 0xE8, 0xE8)
DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
GREEN_TEXT = RGBColor(0x2E, 0x7D, 0x32)


# ── Figure helpers ────────────────────────────────────────────────────────────
def _fig_to_stream(fig) -> BytesIO:
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150, facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_picture(slide, fig, left, top, width, height):
    stream = _fig_to_stream(fig)
    slide.shapes.add_picture(stream, left, top, width, height)


# ── Slide/shape helpers ───────────────────────────────────────────────────────
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_header(slide, title, palette):
    shp = slide.shapes.add_shape(RECT, 0, 0, SW, HDR)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(palette["primary"])
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.size = Pt(22)


def _add_insight_box(slide, heading, body, left, top, width, height, palette):
    shp = slide.shapes.add_shape(RECT, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = _rgb(palette["primary"])
    shp.line.width = Pt(1.5)

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top   = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = heading
    run.font.bold = True
    run.font.size = Pt(10.5)
    run.font.color.rgb = _rgb(palette["primary"])

    for line in (body or "").split("\n"):
        line = line.strip()
        para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(1)
        r = para.add_run()
        r.text = line if line else ""
        r.font.size = Pt(9)
        r.font.color.rgb = DARK_TEXT


# ── Table cell helper ─────────────────────────────────────────────────────────
def _set_cell(cell, text, size=9, bold=False, color=None,
              align=PP_ALIGN.CENTER, bg=None):
    cell.text = str(text)
    tf = cell.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK_TEXT
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


# ── Table builders ────────────────────────────────────────────────────────────
def _summary_table(slide, display_df, date_label, left, top, width, palette):
    primary = _rgb(palette["primary"])
    ROW_H = Inches(0.38)
    n = 5
    tbl = slide.shapes.add_table(n, 4, left, top, width, ROW_H * n).table
    tbl.columns[0].width = Inches(1.7)
    for i in range(1, 4):
        tbl.columns[i].width = Inches(1.1)

    headers = [date_label, "Impressions", "Clicks", "CTR"]
    for c, h in enumerate(headers):
        _set_cell(tbl.cell(0, c), h, bold=True, color=WHITE, bg=primary)

    labels = ["Target", "Delivery", "Over Delivery", "% of Over Delivery"]
    for r, label in enumerate(labels, 1):
        row = display_df.iloc[r - 1]
        bg = LIGHT_RED if label == "Over Delivery" else (LIGHT_GRAY if r % 2 == 0 else WHITE)
        bold = label in ("Delivery", "Over Delivery")
        txt_color = GREEN_TEXT if label == "% of Over Delivery" else DARK_TEXT
        _set_cell(tbl.cell(r, 0), label, bold=bold, color=txt_color, align=PP_ALIGN.LEFT, bg=bg)
        for c, col in enumerate(["Impressions", "Clicks", "CTR"], 1):
            _set_cell(tbl.cell(r, c), row[col], bold=bold, color=txt_color, bg=bg)


def _generic_table(slide, rows_data, headers, col_widths, left, top, width, palette):
    """rows_data: list of lists (no header row). Returns table height."""
    primary = _rgb(palette["primary"])
    ROW_H = Inches(0.35)
    n = 1 + len(rows_data)
    tbl = slide.shapes.add_table(n, len(headers), left, top, width, ROW_H * n).table
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for c, h in enumerate(headers):
        _set_cell(tbl.cell(0, c), h, bold=True, color=WHITE, bg=primary)

    for r, row in enumerate(rows_data, 1):
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        bold = (r == 1)
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            _set_cell(tbl.cell(r, c), val, size=8, bold=bold,
                      align=align, bg=bg)

    return ROW_H * n


# ── Individual slides ─────────────────────────────────────────────────────────
def _title_slide(prs, client, campaign_name, palette, logo_bytes):
    slide = _blank_slide(prs)

    left_rect = slide.shapes.add_shape(RECT, 0, 0, SW / 2, SH)
    left_rect.fill.solid()
    left_rect.fill.fore_color.rgb = WHITE
    left_rect.line.fill.background()

    if logo_bytes:
        logo_w, logo_h = Inches(4.0), Inches(2.5)
        slide.shapes.add_picture(
            BytesIO(logo_bytes),
            (SW / 2 - logo_w) / 2,
            (SH - logo_h) / 2,
            logo_w, logo_h,
        )
    else:
        txb = slide.shapes.add_textbox(PAD, SH / 2 - Inches(0.75), SW / 2 - PAD * 2, Inches(1.5))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = client
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = _rgb(palette["primary"])

    right_rect = slide.shapes.add_shape(RECT, SW / 2, 0, SW / 2, SH)
    right_rect.fill.solid()
    right_rect.fill.fore_color.rgb = _rgb(palette["primary"])
    right_rect.line.fill.background()

    txb = slide.shapes.add_textbox(
        SW / 2 + PAD, SH / 2 - Inches(1.2),
        SW / 2 - PAD * 2, Inches(2.4),
    )
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = campaign_name
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _summary_slide(prs, dv360_data, target_impressions, target_clicks,
                   target_ctr, highlights, campaign_name, palette):
    slide = _blank_slide(prs)
    _add_header(slide, "Weekly Campaign Performance Summary.", palette)

    imp   = int(dv360_data["Impressions"].sum())
    clk   = int(dv360_data["Clicks"].sum())
    ctr   = clk / imp if imp else 0
    o_imp = imp - target_impressions
    o_clk = clk - target_clicks
    o_ctr = ctr - target_ctr
    p_imp = o_imp / target_impressions if target_impressions else 0
    p_clk = o_clk / target_clicks if target_clicks else 0
    p_ctr = o_ctr / target_ctr if target_ctr else 0

    display = pd.DataFrame({
        "Impressions": [f"{int(target_impressions):,}", f"{imp:,}", f"{int(o_imp):,}", f"{p_imp:.2%}"],
        "Clicks":      [f"{int(target_clicks):,}",      f"{clk:,}", f"{int(o_clk):,}", f"{p_clk:.2%}"],
        "CTR":         [f"{target_ctr:.2%}",            f"{ctr:.2%}", f"{o_ctr:.2%}", f"{p_ctr:.2%}"],
    }, index=["Target", "Delivery", "Over Delivery", "% of Over Delivery"])

    content_top = HDR + PAD
    left_w  = Inches(5.4)
    r_start = Inches(5.8)
    r_w     = SW - r_start - PAD

    dates = pd.to_datetime(dv360_data["Date"])
    date_label = (dates.min().strftime("%b %d")
                  + " – " + dates.max().strftime("%b %d"))

    _add_insight_box(slide, f"Campaign Highlights ({date_label})", highlights,
                     PAD, content_top, left_w - PAD, Inches(2.7), palette)

    _summary_table(slide, display, date_label,
                   PAD, content_top + Inches(2.85), left_w - PAD, palette)

    fig = overview_chart(imp, clk, ctr, date_label, palette)
    _add_picture(slide, fig, r_start, content_top, r_w, SH - content_top - PAD)


def _demographics_slide(prs, dv360_data, insights, campaign_name, palette):
    slide = _blank_slide(prs)
    _add_header(slide, "Summary – Demographics Level.", palette)

    demo = (dv360_data.groupby("Demographics")
            .agg(Impressions=("Impressions", "sum"), Clicks=("Clicks", "sum")))
    demo["CTR"] = demo["Clicks"] / demo["Impressions"]

    content_top = HDR + PAD
    left_w  = Inches(5.4)
    r_start = Inches(5.8)
    r_w     = SW - r_start - PAD

    dates = pd.to_datetime(dv360_data["Date"])
    date_label = (dates.min().strftime("%b %d")
                  + " – " + dates.max().strftime("%b %d"))

    rows = []
    for name, row in demo.iterrows():
        rows.append([name, f"{int(row['Impressions']):,}", f"{int(row['Clicks']):,}", f"{row['CTR']:.2%}"])
    total_imp = int(demo["Impressions"].sum())
    total_clk = int(demo["Clicks"].sum())
    rows.append(["Total", f"{total_imp:,}", f"{total_clk:,}",
                 f"{total_clk / total_imp:.2%}" if total_imp else "0.00%"])

    _add_insight_box(slide, f"Demographics Performance ({date_label})", insights,
                     PAD, content_top, left_w - PAD, Inches(2.7), palette)

    _generic_table(slide, rows, [campaign_name, "Impressions", "Clicks", "CTR"],
                   [1.5, 1.15, 1.0, 1.0],
                   PAD, content_top + Inches(2.85), left_w - PAD, palette)

    fig = demographic_chart(demo, palette)
    _add_picture(slide, fig, r_start, content_top, r_w, SH - content_top - PAD)


def _area_slide(prs, dv360_data, insights, campaign_name, palette):
    slide = _blank_slide(prs)
    _add_header(slide, "Summary – Area Level.", palette)

    area_df = (dv360_data.groupby("Store")
               .agg(Impressions=("Impressions", "sum"), Clicks=("Clicks", "sum")))
    area_df["CTR"] = area_df["Clicks"] / area_df["Impressions"]
    area_sorted = area_df.sort_values("Impressions", ascending=False)

    content_top = HDR + PAD
    chart_w = Inches(7.2)
    r_start = Inches(7.5)
    r_w     = SW - r_start - PAD

    dates = pd.to_datetime(dv360_data["Date"])
    date_label = (dates.min().strftime("%b %d")
                  + " – " + dates.max().strftime("%b %d"))

    fig = area_chart(area_df, palette)
    _add_picture(slide, fig, PAD, content_top, chart_w - PAD, SH - content_top - PAD)

    rows = []
    for name, row in area_sorted.iterrows():
        rows.append([name, f"{int(row['Impressions']):,}", f"{int(row['Clicks']):,}", f"{row['CTR']:.2%}"])

    _add_insight_box(slide, f"Area Performance ({date_label})", insights,
                     r_start, content_top, r_w, Inches(2.7), palette)

    _generic_table(slide, rows, ["Area", "Impressions", "Clicks", "CTR"],
                   [1.4, 1.1, 0.9, 0.9],
                   r_start, content_top + Inches(2.85), r_w, palette)


def _products_slide(prs, ft_data, insights, campaign_name, palette):
    slide = _blank_slide(prs)
    _add_header(slide, "Summary – Products.", palette)

    col = (
        "Product" if "Product" in ft_data.columns
        else "Products" if "Products" in ft_data.columns
        else None
    )
    if col is None:
        # No product mapping available — skip this slide
        return
    products = (ft_data.groupby(col)["Clicks"]
                .sum().sort_values(ascending=False).head(10))

    content_top = HDR + PAD
    left_w  = Inches(5.0)
    r_start = Inches(5.3)
    r_w     = SW - r_start - PAD
    insight_h = Inches(1.8)
    insight_top = SH - PAD - insight_h

    rows = [[name, f"{int(clk):,}"] for name, clk in products.items()]
    _generic_table(slide, rows, ["Product Name", "Clicks"],
                   [3.5, 0.9],
                   PAD, content_top, left_w - PAD, palette)

    fig = product_pie_chart(products, palette)
    _add_picture(slide, fig, r_start, content_top, r_w, insight_top - content_top - PAD * 0.5)

    _add_insight_box(slide, f"Product Click Performance ({campaign_name})", insights,
                     PAD, insight_top, SW - PAD * 2, insight_h, palette)


def _creatives_slide(prs, dv360_data, insights, campaign_name, palette):
    slide = _blank_slide(prs)
    _add_header(slide, "Summary – Creatives.", palette)

    creative = (dv360_data.groupby("Creative Size")
                .agg(Impressions=("Impressions", "sum"), Clicks=("Clicks", "sum")))
    creative["CTR"] = creative["Clicks"] / creative["Impressions"]
    creative_sorted = creative.sort_values("Clicks", ascending=False)

    content_top = HDR + PAD
    left_w  = Inches(5.4)
    r_start = Inches(5.8)
    r_w     = SW - r_start - PAD

    rows = []
    for name, row in creative_sorted.iterrows():
        rows.append([name, f"{int(row['Impressions']):,}", f"{int(row['Clicks']):,}", f"{row['CTR']:.2%}"])

    tbl_h = _generic_table(slide, rows, ["Creative Size", "Impressions", "Clicks", "CTR"],
                           [1.3, 1.1, 0.9, 0.9],
                           PAD, content_top, left_w - PAD, palette)

    insight_top = content_top + tbl_h + Inches(0.2)
    insight_h   = SH - insight_top - PAD
    dates = pd.to_datetime(dv360_data["Date"])
    date_label  = (dates.min().strftime("%b %d")
                   + " – " + dates.max().strftime("%b %d"))
    _add_insight_box(slide, f"Creative Performance ({date_label})", insights,
                     PAD, insight_top, left_w - PAD, max(insight_h, Inches(1.0)), palette)

    fig = creative_chart(creative, palette)
    _add_picture(slide, fig, r_start, content_top, r_w, SH - content_top - PAD)


# ── Public entry point ────────────────────────────────────────────────────────
def build_pptx(
    client,
    campaign_name,
    habanero_df,
    ft_data,
    dv360_data,
    target_impressions,
    target_clicks,
    target_ctr,
    highlights,
    demographic_insights,
    area_insights,
    product_insights,
    creative_insights,
    palette,
    logo_bytes=None,
) -> BytesIO:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    _title_slide(prs, client, campaign_name, palette, logo_bytes)
    _summary_slide(prs, dv360_data, target_impressions, target_clicks,
                   target_ctr, highlights, campaign_name, palette)
    _demographics_slide(prs, dv360_data, demographic_insights, campaign_name, palette)
    _area_slide(prs, dv360_data, area_insights, campaign_name, palette)
    _products_slide(prs, ft_data, product_insights, campaign_name, palette)
    _creatives_slide(prs, dv360_data, creative_insights, campaign_name, palette)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
